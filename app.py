import os
import json
import time
from io import BytesIO
from PIL import Image

# PowerPoint Libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Google GenAI SDK
from google import genai
from google.genai import types

# --- CONFIGURATION ---
# It is best practice to use environment variables for Project IDs
PROJECT_ID = os.getenv("GOOGLE_CLOUD_PROJECT_ID", "your-gcp-project-id")
LOCATION = "us-central1"

# Initialize the client
# Ensure you have run 'gcloud auth application-default login' in your terminal
client = genai.Client(vertexai=True, project=PROJECT_ID, location=LOCATION)

# Define the models
TEXT_MODEL_NAME = "gemini-2.5-pro"
IMAGE_MODEL_NAME = "gemini-2.0-flash-preview-image-generation"

def generate_presentation_slide_text(topic, num_slides=5):
    """Generates a presentation outline in JSON format."""
    outline_prompt = f"""
    Create a detailed, structured presentation outline in JSON format for the topic: "{topic}".
    The JSON object should have a "title" field for the main presentation title, and a "slides" array.
    Create exactly {num_slides} slides.
    Each object in the "slides" array must contain three fields:
    1. "slide_title": A concise title for the slide.
    2. "key_points": An array of 3-5 bullet points summarizing the content.
    3. "image_prompt": A descriptive, single-sentence prompt to generate a relevant image for this slide.
    """
    print(f"--- Generating outline for: {topic} ---")

    try:
        outline_response = client.models.generate_content(
            model=TEXT_MODEL_NAME,
            contents=outline_prompt
        )
        
        # Clean the response text to ensure it's valid JSON
        json_text = outline_response.candidates[0].content.parts[0].text.strip('`').strip('json').strip()
        return json.loads(json_text)
    except Exception as e:
        print(f"Failed to generate or parse outline: {e}")
        return None

def generate_presentation_slide_images(outline):
    """Generates images for each slide based on the outline."""
    images = {}
    print("\n--- Generating images for slides ---")

    for i, slide in enumerate(outline['slides']):
        image_prompt = slide['image_prompt']
        print(f"Processing Slide {i+1}: {slide['slide_title']}")
        try:
            response = client.models.generate_content(
                model=IMAGE_MODEL_NAME,
                contents=image_prompt,
                config=types.GenerateContentConfig(
                    response_modalities=['TEXT', 'IMAGE']
                )
            )

            for part in response.candidates[0].content.parts:
                if part.inline_data is not None:
                    images[i] = part.inline_data.data
                    break
            else:
                images[i] = None
        except Exception as e:
            print(f"Error generating image for slide {i+1}: {e}")
            images[i] = None

    return images

def create_presentation_file(outline, images, output_filename="AI_Presentation.pptx"):
    """Assembles the PowerPoint presentation."""
    prs = Presentation()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = outline['title']
    
    # Simple styling for title
    title_p = title_shape.text_frame.paragraphs[0]
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0, 51, 102)

    subtitle = slide.placeholders[1]
    subtitle.text = f"AI-Generated Presentation: {outline['title']}"

    # 2. Content Slides
    content_slide_layout = prs.slide_layouts[1]
    for i, slide_data in enumerate(outline['slides']):
        slide = prs.slides.add_slide(content_slide_layout)

        # Clear default shapes
        for shape in list(slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)

        # Background Header
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0, 51, 102)
        rect.line.fill.background()

        # Slide Title
        title_textbox = slide.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, Inches(1))
        title_frame = title_textbox.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data['slide_title']
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.font.size = Pt(24)
        title_p.alignment = PP_ALIGN.CENTER
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Bullet Points (Left Side)
        body_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0))
        text_frame = body_textbox.text_frame
        text_frame.word_wrap = True

        for point in slide_data['key_points']:
            p = text_frame.add_paragraph()
            p.text = '• ' + point
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 30, 60)
            p.space_after = Pt(10)

        # Image (Right Side)
        if i in images and images[i] is not None:
            image_stream = BytesIO(images[i])
            slide.shapes.add_picture(image_stream, Inches(5.5), Inches(2.0), width=Inches(4.0))

    prs.save(output_filename)
    print(f"\nSuccess! File saved as: {output_filename}")

# --- MAIN EXECUTION ---
if __name__ == "__main__":
    topic = "The Future of Quantum Computing"
    num_slides = 5
    
    presentation_outline = generate_presentation_slide_text(topic, num_slides)
    
    if presentation_outline:
        slide_images = generate_presentation_slide_images(presentation_outline)
        create_presentation_file(presentation_outline, slide_images)
